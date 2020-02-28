from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

def rectanglepptx(shapes,offset,sx,sy,sw,sh,text):
    #print(f"x:{sx:.2f}|y:{sy-sh:.2f}|w:{sw:.2f}|h:{sh:.2f}")
    isx = Inches((sx+offset)/2.54)
    isy = Inches((sy-sh)/2.54)
    isw = Inches(sw/2.54)
    ish = Inches(sh/2.54)
    # isx,isy is the top left corner of the rectangle
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE,isx,isy,isw,ish)
    shape.shadow.inherit = False
    shape.fore_color = RGBColor(96, 96, 96)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    line = shape.line
    line.color.rgb = RGBColor(96, 96, 96)
    text_frame = shape.text_frame
    text_frame.margin_left=Inches(0.0) 
    text_frame.margin_right=Inches(0.0) 
    text_frame.margin_top=Inches(0.0)
    text_frame.margin_bottom=Inches(0.0)
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE                
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(8)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = RGBColor(96, 96, 96)
    return (run,font,fill)
    
def marimekkopptx(ssmarkets,ccsss,lowshare,drawshare,marketshare,height,width,showther,marketsize,name):
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = f"Market Size = {float(marketsize):,.0f}"
    
    x0 = 1.0
    y0 = 17.0
    owidth = 3.0

    x=0
    y=0
    marketcount=0

    for market in ssmarkets.index:
        h = float(ssmarkets.loc[market])
        #print(f"height'{market}':{h:.0f}")
        competitors = ccsss.loc[market].sort_values(by=['Share'],ascending=False)
        x=0
        other = 0.0
        count=0
        first = True
        ccompetitor=0
        for competitor in competitors.index:
            if ccompetitor == 0:
                # print market label
                sx = (x*width+x0)
                sy = (y0-(y*height))
                sw = owidth
                sh = (h*height)
                run,font,fill = rectanglepptx(shapes,0,sx,sy,sw,sh,f"{market} ({h*100:.0f}%)")
                fill.fore_color.rgb = RGBColor(192, 192, 192)
                font.color.rgb = RGBColor(255, 255, 255)
            if competitor != 'Other':
                w = float(ccsss.loc[market].loc[competitor])
                if (w >=lowshare):
                    #print(f"  width {competitor}:{w:.0f}")
                    #print(f"({w},{h})")
                    # create the rectangle
                    sx = (x*width+x0)
                    sy = (y0-(y*height))
                    sw = (w*width)
                    sh = (h*height)
                    run,font,fill = rectanglepptx(shapes,owidth,sx,sy,sw,sh,"")
                    label = competitor.lower()
                    if name.lower() in label:
                        fill.fore_color.rgb = RGBColor(0, 204, 102)
                    if w<=drawshare:
                        run.text = f"{competitor.replace('Machine: ','')[:5]} ({w*100:.0f}%)"
                        #run.text = f"{competitor} ({w*100:.0f}%)"
                        if h >= marketshare:
                            font.size = Pt(6)
                        else:
                            font.size = Pt(4)
                    else:
                        run.text = f"{competitor} ({w*100:.0f}%)"
                        if h >= marketshare:
                            font.size = Pt(8)
                        else:
                            font.size = Pt(6)
                    x += w
                else:
                    other += w
            else:
                other += float(ccsss.loc[market].loc['Other'])
            ccompetitor +=1
                
        #print(x,y)
        w=other
        if w >0:
            sx = (x*width+x0)
            sy = (y0-(y*height))
            sw = (w*width)
            sh = (h*height)
            run,font,fill = rectanglepptx(shapes,owidth,sx,sy,sw,sh,"")
            if w<=drawshare:
                run.text = f"({w*100:.0f}%)"
                if h >= marketshare:
                    font.size = Pt(6)
                else:
                    font.size = Pt(4)
            else:
                run.text = f"({w*100:.0f}%)"
                if h >= marketshare:
                    font.size = Pt(8)
                else:
                    font.size = Pt(6)
        y += h
        marketcount += 1

    return prs